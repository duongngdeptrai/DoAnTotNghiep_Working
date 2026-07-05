# -*- coding: utf-8 -*-
"""Dựng slide thuyết trình bảo vệ đồ án tốt nghiệp từ nội dung LaTeX."""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from PIL import Image

HERE = os.path.dirname(os.path.abspath(__file__))
IMG_DIR = os.path.join(HERE, "..", "SOICT_DATN_Duong", "Hinhve")
OUT_PATH = os.path.join(HERE, "DATN_NguyenDucDuong_SlideBaoVe.pptx")

def IMG(name):
    return os.path.join(IMG_DIR, name)

# ---------------------------------------------------------------- THEME ----
# Do HUST (#C4151C) + vang gold (#FFCB09) CHI dung cho tieu de slide / bia /
# slide cam on (header(), cover, closing). Toan bo phan noi dung ben duoi
# tieu de dung mau trung tinh (den/xam) — khong con "mau me" o than slide.
NAVY      = RGBColor(0x7A, 0x0D, 0x12)   # đỏ đậm HUST — CHỈ dùng cho bìa/slide cảm ơn
NAVY_2    = RGBColor(0x9E, 0x1D, 0x22)   # đỏ đậm phụ — CHỈ dùng cho bìa/slide cảm ơn
GOLD      = RGBColor(0xFF, 0xCB, 0x09)   # vàng gold tươi — CHỈ dùng cho gạch chân tiêu đề/bìa
HEADER_RED = RGBColor(0xC4, 0x15, 0x1C)  # đỏ HUST — CHỈ dùng cho thanh trên cùng + kicker của header()

# Mau trung tinh dung cho toan bo noi dung (card, bullet, bang, callout...)
TEAL      = RGBColor(0x33, 0x33, 0x33)   # xám đậm — accent chính trong nội dung (thay đỏ)
TEAL_LT   = RGBColor(0xF4, 0xF4, 0xF4)   # nền card nhạt (xám rất nhạt)
CORAL     = RGBColor(0x59, 0x59, 0x59)   # xám trung — accent phụ trong nội dung (thay vàng)
CORAL_LT  = RGBColor(0xEA, 0xEA, 0xEA)   # nền card nhạt thứ cấp (xám nhạt hơn 1 chút)
DARK_BOX  = RGBColor(0x2B, 0x2B, 0x2B)   # khối nền tối trung tính trong nội dung (thay NAVY)
DARK_BOX_2= RGBColor(0x50, 0x50, 0x50)   # khối nền tối phụ trong nội dung (thay NAVY_2)
BG        = RGBColor(0xFF, 0xFF, 0xFF)
BG_SOFT   = RGBColor(0xF6, 0xF8, 0xFA)
TEXT_DARK = RGBColor(0x22, 0x22, 0x22)
TEXT_MUTE = RGBColor(0x59, 0x59, 0x59)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
LINE_GRAY = RGBColor(0xDD, 0xDD, 0xDD)

FONT = "Calibri"

ASSET_DIR = os.path.join(HERE, "assets")
LOGO = os.path.join(ASSET_DIR, "hust_logo.png")
COVER_BG = os.path.join(ASSET_DIR, "hust_bg_cover.jpg")
SIDEBAR_BG = os.path.join(ASSET_DIR, "hust_bg_sidebar.jpg")
LOGO_AR = 300 / 390  # tỷ lệ rộng/cao thật của hust_logo.png

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
BLANK = prs.slide_layouts[6]

SECTION_TAGS = {
    "problem": "01  ĐẶT VẤN ĐỀ",
    "design":  "02  KIẾN TRÚC & THIẾT KẾ",
    "tech":    "03  CÔNG NGHỆ SỬ DỤNG",
    "contrib": "04  ĐÓNG GÓP KỸ THUẬT NỔI BẬT",
    "result":  "05  KẾT QUẢ & ĐÁNH GIÁ",
    "concl":   "06  KẾT LUẬN",
}

PAGE_COUNTER = {"n": 0}

# ------------------------------------------------------------- HELPERS -----

def new_slide():
    return prs.slides.add_slide(BLANK)

def set_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def _no_line(shape):
    shape.line.fill.background()

def rect(slide, l, t, w, h, color, line=False):
    sp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    sp.fill.solid()
    sp.fill.fore_color.rgb = color
    if line:
        sp.line.color.rgb = LINE_GRAY
        sp.line.width = Pt(0.75)
    else:
        _no_line(sp)
    sp.shadow.inherit = False
    return sp

def rounded_rect(slide, l, t, w, h, color, radius=0.06, line=False, line_color=None):
    sp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    try:
        sp.adjustments[0] = radius
    except Exception:
        pass
    sp.fill.solid()
    sp.fill.fore_color.rgb = color
    if line:
        sp.line.color.rgb = line_color or LINE_GRAY
        sp.line.width = Pt(1)
    else:
        _no_line(sp)
    sp.shadow.inherit = False
    return sp

def text(slide, l, t, w, h, s, size=18, color=TEXT_DARK, bold=False, italic=False,
          align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font=FONT, line_spacing=1.0,
          letter_spacing=None, wrap=True):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    lines = s.split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        r = p.add_run()
        r.text = line
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.italic = italic
        r.font.name = font
        r.font.color.rgb = color
    return tb

def bullets(slide, l, t, w, h, items, size=16, color=TEXT_DARK, bullet_color=TEAL,
            space_after=8, bold_lead=False, font=FONT, anchor=MSO_ANCHOR.TOP):
    """items: list of str, or (str, level) tuples."""
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    for i, item in enumerate(items):
        level = 0
        if isinstance(item, tuple):
            item_text, level = item
        else:
            item_text = item
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(space_after)
        p.line_spacing = 1.08
        indent = Inches(0.28 * level)
        p.level = 0
        pPr = p._pPr
        if pPr is None:
            pPr = p._p.get_or_add_pPr()
        pPr.set('marL', str(indent))
        pPr.set('indent', str(-Inches(0.22)))
        # bullet char + color
        buClr = pPr.makeelement(qn('a:buClr'))
        srgb = buClr.makeelement(qn('a:srgbClr'), {'val': str(bullet_color)})
        buClr.append(srgb)
        buFont = pPr.makeelement(qn('a:buFont'), {'typeface': font})
        buChar = pPr.makeelement(qn('a:buChar'), {'char': '—' if level == 0 else '·'})
        pPr.append(buClr)
        pPr.append(buFont)
        pPr.append(buChar)
        r = p.add_run()
        r.text = item_text
        r.font.size = Pt(size - level * 1)
        r.font.name = font
        r.font.color.rgb = color
        r.font.bold = bold_lead and level == 0
    return tb

def numbered_steps(slide, l, t, w, items, size=13.5, color=TEXT_DARK, badge_color=TEAL,
                     item_h=0.52, gap=0.1):
    """Danh sách bước có huy hiệu đánh số tròn — dùng để trình bày luồng quy trình."""
    y = t
    badge_d = Inches(0.34)
    for i, item in enumerate(items, start=1):
        circ = slide.shapes.add_shape(MSO_SHAPE.OVAL, l, y, badge_d, badge_d)
        circ.fill.solid()
        circ.fill.fore_color.rgb = badge_color
        _no_line(circ)
        circ.shadow.inherit = False
        tf = circ.text_frame
        tf.margin_left = 0
        tf.margin_right = 0
        tf.margin_top = 0
        tf.margin_bottom = 0
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = str(i)
        r.font.size = Pt(13)
        r.font.bold = True
        r.font.color.rgb = WHITE
        r.font.name = FONT
        text(slide, l + badge_d + Inches(0.15), y - Inches(0.02), w - badge_d - Inches(0.15),
             Inches(item_h), item, size=size, color=color, line_spacing=1.15)
        y += Inches(item_h) + Inches(gap)
    return y

def header(slide, tag, title, subtitle=None):
    """Chuẩn header cho slide nội dung: kicker + tiêu đề + gạch chân accent + logo HUST."""
    rect(slide, 0, 0, SLIDE_W, Inches(0.09), HEADER_RED)
    logo_h = Inches(0.42)
    logo_w = Emu(int(logo_h * LOGO_AR))
    slide.shapes.add_picture(LOGO, SLIDE_W - Inches(0.55) - logo_w, Inches(0.24), height=logo_h, width=logo_w)
    text(slide, Inches(0.55), Inches(0.32), Inches(9), Inches(0.3), tag,
         size=12, color=HEADER_RED, bold=True, font=FONT)
    text(slide, Inches(0.55), Inches(0.6), Inches(11), Inches(0.75), title,
         size=30, color=HEADER_RED, bold=True, font=FONT)
    rect(slide, Inches(0.57), Inches(1.32), Inches(0.6), Inches(0.05), GOLD)
    if subtitle:
        text(slide, Inches(0.55), Inches(1.42), Inches(11.8), Inches(0.4), subtitle,
             size=14, color=TEXT_MUTE, italic=True)

def footer(slide, label="HỆ THỐNG THEO DÕI VỊ TRÍ TRẺ EM THEO THỜI GIAN THỰC"):
    PAGE_COUNTER["n"] += 1
    rect(slide, 0, Inches(7.34), SLIDE_W, Pt(0.75), LINE_GRAY)
    text(slide, Inches(0.55), Inches(7.36), Inches(8), Inches(0.3), label,
         size=9, color=TEXT_MUTE)
    text(slide, Inches(12.4), Inches(7.36), Inches(0.5), Inches(0.3), str(PAGE_COUNTER["n"]),
         size=9, color=TEXT_MUTE, align=PP_ALIGN.RIGHT)

def content_slide(tag, title, subtitle=None, bg=BG):
    s = new_slide()
    set_bg(s, bg)
    header(s, tag, title, subtitle)
    footer(s)
    return s

def image_fit(slide, path, l, t, max_w, max_h, align="center", valign="middle", border=True):
    im = Image.open(path)
    iw, ih = im.size
    ar = iw / ih
    box_ar = max_w / max_h
    if ar > box_ar:
        w = max_w
        h = w / ar
    else:
        h = max_h
        w = h * ar
    if align == "center":
        x = l + (max_w - w) / 2
    elif align == "left":
        x = l
    else:
        x = l + (max_w - w)
    if valign == "middle":
        y = t + (max_h - h) / 2
    elif valign == "top":
        y = t
    else:
        y = t + (max_h - h)
    if border:
        pad = Emu(int(Pt(3)))
        rect(slide, x - pad, y - pad, w + 2 * pad, h + 2 * pad, WHITE, line=True)
    pic = slide.shapes.add_picture(path, x, y, width=Emu(int(w)), height=Emu(int(h)))
    return pic

def card(slide, l, t, w, h, title_txt, body_items, title_color=TEAL, bg=TEAL_LT,
          title_size=16, body_size=13.5):
    rounded_rect(slide, l, t, w, h, bg, radius=0.045)
    rect(slide, l, t, Inches(0.07), h, title_color)
    text(slide, l + Inches(0.3), t + Inches(0.18), w - Inches(0.5), Inches(0.4),
         title_txt, size=title_size, bold=True, color=TEXT_DARK)
    bullets(slide, l + Inches(0.3), t + Inches(0.65), w - Inches(0.55), h - Inches(0.8),
            body_items, size=body_size, bullet_color=title_color, space_after=5)

def metric_card(slide, l, t, w, h, value, label, color=WHITE):
    rounded_rect(slide, l, t, w, h, DARK_BOX_2, radius=0.08)
    text(slide, l, t + Inches(0.18), w, Inches(0.75), value, size=30, bold=True,
         color=color, align=PP_ALIGN.CENTER)
    text(slide, l + Inches(0.15), t + h - Inches(0.55), w - Inches(0.3), Inches(0.45),
         label, size=12, color=WHITE, align=PP_ALIGN.CENTER, line_spacing=1.0)

def make_table(slide, l, t, w, h, data, col_w_ratio=None, header_bg=DARK_BOX,
                header_color=WHITE, body_size=13, header_size=13.5, row_h=None,
                highlight_col=None, highlight_color=TEAL_LT):
    rows = len(data)
    cols = len(data[0])
    gtable = slide.shapes.add_table(rows, cols, l, t, w, h).table
    if col_w_ratio:
        total = sum(col_w_ratio)
        for i, ratio in enumerate(col_w_ratio):
            gtable.columns[i].width = Emu(int(w * ratio / total))
    if row_h:
        for r in range(rows):
            gtable.rows[r].height = row_h
    for c in range(cols):
        cell = gtable.cell(0, c)
        cell.fill.solid()
        cell.fill.fore_color.rgb = header_bg
        cell.margin_left = Inches(0.08)
        cell.margin_right = Inches(0.08)
        cell.margin_top = Inches(0.04)
        cell.margin_bottom = Inches(0.04)
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf = cell.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = data[0][c]
        r.font.size = Pt(header_size)
        r.font.bold = True
        r.font.color.rgb = header_color
        r.font.name = FONT
    for ri in range(1, rows):
        for c in range(cols):
            cell = gtable.cell(ri, c)
            cell.fill.solid()
            is_hl = (highlight_col is not None and c == highlight_col)
            if is_hl:
                cell.fill.fore_color.rgb = highlight_color
            else:
                cell.fill.fore_color.rgb = WHITE if ri % 2 else BG_SOFT
            cell.margin_left = Inches(0.08)
            cell.margin_right = Inches(0.08)
            cell.margin_top = Inches(0.03)
            cell.margin_bottom = Inches(0.03)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            tf = cell.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT if c == 0 else PP_ALIGN.CENTER
            r = p.add_run()
            r.text = data[ri][c]
            r.font.size = Pt(body_size)
            r.font.bold = is_hl
            r.font.color.rgb = TEAL if is_hl else TEXT_DARK
            r.font.name = FONT
    # remove default banding style (first-row/band formatting) is fine to leave.
    return gtable

# =====================================================================
# SLIDE 0 — BÌA
# =====================================================================
s = new_slide()
set_bg(s, NAVY)
rect(s, 0, 0, SLIDE_W, Inches(0.12), GOLD)
rect(s, 0, Inches(7.38), SLIDE_W, Inches(0.12), GOLD)

cover_logo_h = Inches(0.9)
cover_logo_w = Emu(int(cover_logo_h * LOGO_AR))
s.shapes.add_picture(LOGO, (SLIDE_W - cover_logo_w) / 2, Inches(0.35), height=cover_logo_h, width=cover_logo_w)
text(s, Inches(1), Inches(1.32), Inches(11.3), Inches(0.35), "TRƯỜNG ĐẠI HỌC BÁCH KHOA HÀ NỘI",
     size=14, color=GOLD, bold=True, align=PP_ALIGN.CENTER)
text(s, Inches(1), Inches(1.63), Inches(11.3), Inches(0.3), "HANOI UNIVERSITY OF SCIENCE AND TECHNOLOGY",
     size=10.5, color=RGBColor(0xB9, 0xCB, 0xD8), align=PP_ALIGN.CENTER)

text(s, Inches(1), Inches(2.12), Inches(11.3), Inches(0.4), "ĐỒ ÁN TỐT NGHIỆP",
     size=16, color=GOLD, bold=True, align=PP_ALIGN.CENTER, letter_spacing=True)
text(s, Inches(1), Inches(2.62), Inches(11.3), Inches(1.6),
     "HỆ THỐNG THEO DÕI VỊ TRÍ TRẺ EM\nTHEO THỜI GIAN THỰC",
     size=34, color=WHITE, bold=True, align=PP_ALIGN.CENTER, line_spacing=1.1)
rect(s, Inches(5.67), Inches(4.15), Inches(2), Inches(0.035), GOLD)
text(s, Inches(1.5), Inches(4.32), Inches(10.3), Inches(0.5),
     "Thiết bị IoT (ESP32 + SIM7600)  ·  Backend FastAPI + MongoDB  ·  Ứng dụng di động Flutter",
     size=15, color=RGBColor(0xB9, 0xCB, 0xD8), align=PP_ALIGN.CENTER)

rounded_rect(s, Inches(3.67), Inches(4.85), Inches(6), Inches(1.55), NAVY_2, radius=0.12)
text(s, Inches(3.67), Inches(5.02), Inches(6), Inches(0.35), "Sinh viên thực hiện: Nguyễn Đức Dương",
     size=15, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
text(s, Inches(3.67), Inches(5.36), Inches(6), Inches(0.3), "MSSV: 20225122",
     size=13, color=RGBColor(0xB9, 0xCB, 0xD8), align=PP_ALIGN.CENTER)
text(s, Inches(3.67), Inches(5.72), Inches(6), Inches(0.3), "Giảng viên hướng dẫn: PGS.TS. Tạ Hải Tùng",
     size=13, color=RGBColor(0xB9, 0xCB, 0xD8), align=PP_ALIGN.CENTER)
text(s, Inches(3.67), Inches(6.06), Inches(6), Inches(0.3),
     "Trường Công nghệ Thông tin và Truyền thông – Đại học Bách khoa Hà Nội",
     size=11.5, color=RGBColor(0x8C, 0xA3, 0xB3), align=PP_ALIGN.CENTER)
text(s, Inches(1), Inches(6.9), Inches(11.3), Inches(0.35), "07/2026",
     size=12, color=RGBColor(0x6F, 0x87, 0x99), align=PP_ALIGN.CENTER)

# =====================================================================
# SLIDE 1 — ĐẶT VẤN ĐỀ
# =====================================================================
s = content_slide(SECTION_TAGS["problem"], "Đặt vấn đề")
bullets(s, Inches(0.55), Inches(1.75), Inches(7.6), Inches(5.2), [
    "An toàn trẻ em là mối quan tâm hàng đầu của các bậc phụ huynh trong xã hội hiện đại",
    "Hàng nghìn trường hợp trẻ lạc, mất tích hoặc gặp tai nạn mỗi năm tại các đô thị lớn (Cục Trẻ em – Bộ LĐ-TB&XH)",
    "Trẻ tiểu học / THCS đã di chuyển độc lập nhưng chưa đủ nhận thức để tự bảo vệ bản thân",
    "Phụ huynh thường chỉ phát hiện sự cố khi đã quá muộn để xử lý kịp thời",
], size=17, space_after=16)
rounded_rect(s, Inches(0.55), Inches(5.6), Inches(7.6), Inches(1.2), DARK_BOX, radius=0.08)
text(s, Inches(0.85), Inches(5.75), Inches(7), Inches(0.9),
     "→ Cần theo dõi vị trí thời gian thực, kết hợp cảnh báo tự động\nkhi trẻ rời khỏi vùng an toàn đã định trước",
     size=15.5, color=WHITE, bold=True, line_spacing=1.15)
card(s, Inches(8.5), Inches(1.75), Inches(4.28), Inches(5.05), "Cơ hội công nghệ", [
    "IoT + hạ tầng viễn thông di động phát triển mạnh những năm gần đây",
    "Cho phép xây dựng thiết bị định vị nhỏ gọn, chi phí thấp, hoạt động độc lập",
    "Có thể mở rộng: giám sát người cao tuổi, quản lý học sinh dã ngoại, theo dõi tài sản di động",
], body_size=14, title_size=17)

# =====================================================================
# SLIDE 2 — HẠN CHẾ GIẢI PHÁP HIỆN CÓ
# =====================================================================
s = content_slide(SECTION_TAGS["problem"], "Giải pháp hiện có và hạn chế")
data = [
    ["Tiêu chí", "Life360", "AirTag", "GPS Tracker\nthương mại", "Đề tài này"],
    ["Không cần điện thoại trẻ", "Không", "Có", "Có", "Có"],
    ["Theo dõi thời gian thực", "Có", "Không", "Có", "Có"],
    ["Geofence tự động", "Có", "Không", "Hạn chế", "Có"],
    ["Cảnh báo đa kênh", "Có", "Không", "Hạn chế", "Có"],
    ["Chi phí vận hành", "Miễn phí / thuê bao", "Một lần", "Thuê bao hàng tháng", "Thấp (chỉ SIM data)"],
    ["Khả năng tùy chỉnh", "Thấp", "Rất thấp", "Thấp", "Cao"],
]
make_table(s, Inches(0.55), Inches(1.75), Inches(12.25), Inches(4.55), data,
           col_w_ratio=[2.1, 1.55, 1.35, 1.9, 1.75], highlight_col=4, header_size=14, body_size=13.5)
text(s, Inches(0.55), Inches(6.5), Inches(12.25), Inches(0.7),
     "Hạn chế chung: phụ thuộc điện thoại của trẻ · không theo dõi liên tục thời gian thực · chi phí vận hành cao · khả năng tùy biến thấp",
     size=13.5, color=CORAL, italic=True, bold=True)

# =====================================================================
# SLIDE 3 — MỤC TIÊU & PHẠM VI
# =====================================================================
s = content_slide(SECTION_TAGS["problem"], "Mục tiêu và phạm vi đề tài")
card(s, Inches(0.55), Inches(1.75), Inches(5.95), Inches(3.7), "Mục tiêu", [
    "Thiết bị IoT chuyên dụng hoạt động độc lập, không cần điện thoại của trẻ",
    "Backend xử lý và lưu trữ dữ liệu vị trí thời gian thực",
    "Ứng dụng di động: xem bản đồ, thiết lập geofence, nhận cảnh báo tức thì",
], body_size=15, title_size=18)
card(s, Inches(6.8), Inches(1.75), Inches(5.95), Inches(3.7), "Phạm vi", [
    "Xây dựng nguyên mẫu hoàn chỉnh, có khả năng vận hành thực tế",
    "Bao gồm: thiết kế phần cứng, firmware, backend, ứng dụng di động",
    "Không bao gồm: thiết kế PCB tùy chỉnh, đóng gói sản phẩm thương mại",
], body_size=15, title_size=18, title_color=CORAL, bg=CORAL_LT)
rounded_rect(s, Inches(0.55), Inches(5.75), Inches(12.2), Inches(1.05), DARK_BOX, radius=0.1)
metrics = [("< 5s", "Độ trễ\nend-to-end"), ("~3,2 m", "Độ chính xác\nGPS"), ("5s", "Chu kỳ\ncập nhật")]
mw = Inches(3.8)
gap = Inches(0.3)
startx = Inches(0.55) + (Inches(12.2) - (mw * 3 + gap * 2)) / 2
for i, (val, lab) in enumerate(metrics):
    x = startx + i * (mw + gap)
    text(s, x, Inches(5.85), mw, Inches(0.45), val, size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    text(s, x, Inches(6.3), mw, Inches(0.4), lab.replace("\n", " "), size=12, color=WHITE, align=PP_ALIGN.CENTER)

# =====================================================================
# SLIDE 4 — KIẾN TRÚC TỔNG THỂ
# =====================================================================
s = content_slide(SECTION_TAGS["design"], "Kiến trúc tổng thể hệ thống")
image_fit(s, IMG("kientruc_tongthe.png"), Inches(0.55), Inches(1.75), Inches(12.25), Inches(3.55))
bullets(s, Inches(0.8), Inches(5.55), Inches(11.7), Inches(1.6), [
    "Ba tầng: Thiết bị IoT — Backend xử lý — Ứng dụng di động cho phụ huynh",
    "Kiến trúc hướng sự kiện (event-driven): MQTT (thiết bị → backend) và WebSocket/REST (backend ↔ app)",
    "Mỗi thành phần tổ chức theo kiến trúc phân tầng (layered), phụ thuộc một chiều",
], size=15, space_after=6)

# =====================================================================
# SLIDE 4B — BIỂU ĐỒ GÓI
# =====================================================================
s = content_slide(SECTION_TAGS["design"], "Biểu đồ gói tổng quan hệ thống")
image_fit(s, IMG("bieudo_goi.png"), Inches(0.55), Inches(1.75), Inches(12.25), Inches(4.3))
bullets(s, Inches(0.8), Inches(6.15), Inches(11.7), Inches(1.05), [
    "Ba gói chính: ESP32 Firmware — Backend — Flutter App; gói ngoài: MQTT Broker, MongoDB, Telegram/Email",
    "Tổ chức theo tầng (layered), phụ thuộc một chiều: gói tầng trên phụ thuộc gói tầng dưới, không phụ thuộc ngược hay vượt tầng",
], size=13.5, space_after=5)

# =====================================================================
# SLIDE 5 — CHỨC NĂNG CHÍNH
# =====================================================================
s = content_slide(SECTION_TAGS["design"], "Chức năng chính của hệ thống")
image_fit(s, IMG("usecase_tong_quat.png"), Inches(0.55), Inches(1.75), Inches(5.7), Inches(4.5), border=True)
bullets(s, Inches(6.55), Inches(1.8), Inches(6.2), Inches(4.1), [
    "Đăng ký / Đăng nhập tài khoản",
    "Quản lý thiết bị GPS",
    "Xem vị trí trẻ theo thời gian thực trên bản đồ",
    "Quản lý vùng an toàn địa lý (geofence)",
    "Cảnh báo vi phạm vùng an toàn",
    "Gửi tín hiệu SOS khẩn cấp",
    "Xem thống kê hành trình",
], size=16, space_after=11)
rounded_rect(s, Inches(6.55), Inches(6.35), Inches(6.2), Inches(0.55), TEAL_LT, radius=0.15)
text(s, Inches(6.75), Inches(6.42), Inches(5.9), Inches(0.4),
     "3 tác nhân: Phụ huynh · Thiết bị IoT · Dịch vụ thông báo", size=12.5, bold=True, color=TEAL)
text(s, Inches(0.55), Inches(6.42), Inches(5.7), Inches(0.4),
     "→ 3 luồng nghiệp vụ cốt lõi được trình bày chi tiết ở các trang tiếp theo",
     size=11.5, italic=True, color=TEXT_MUTE)

# =====================================================================
# SLIDE 5B — LUỒNG 1: CẬP NHẬT VỊ TRÍ THỜI GIAN THỰC
# =====================================================================
s = content_slide(SECTION_TAGS["design"], "Luồng nghiệp vụ 1: Cập nhật vị trí thời gian thực",
                   subtitle="Tác nhân: Thiết bị IoT · Backend · Ứng dụng di động")
image_fit(s, IMG("activity_vitri.png"), Inches(0.55), Inches(2.05), Inches(6.0), Inches(4.85), border=True)
numbered_steps(s, Inches(6.9), Inches(2.05), Inches(5.85), [
    "Thiết bị IoT đọc tọa độ GPS mỗi 5 giây (SIM7600)",
    "Gửi bản tin lên MQTT Broker qua kết nối 4G LTE",
    "Backend nhận bản tin, lưu vào MongoDB và kiểm tra geofence tương ứng",
    "Backend đẩy vị trí mới đến ứng dụng đang theo dõi qua WebSocket",
    "Ứng dụng Flutter cập nhật marker trên bản đồ ngay khi nhận được dữ liệu",
], size=14, item_h=0.62, gap=0.14)
rounded_rect(s, Inches(6.9), Inches(6.55), Inches(5.85), Inches(0.65), TEAL_LT, radius=0.12)
text(s, Inches(7.1), Inches(6.63), Inches(5.5), Inches(0.5),
     "Nếu mất tín hiệu GPS/WebSocket: hiển thị vị trí cuối cùng đã biết, tự kết nối lại sau 3s",
     size=11.5, italic=True, color=TEAL, line_spacing=1.1)

# =====================================================================
# SLIDE 5C — LUỒNG 2: VÙNG AN TOÀN & CẢNH BÁO VI PHẠM
# =====================================================================
s = content_slide(SECTION_TAGS["design"], "Luồng nghiệp vụ 2: Vùng an toàn & cảnh báo vi phạm",
                   subtitle="Tác nhân: Phụ huynh · Thiết bị IoT · Backend · Dịch vụ thông báo")
image_fit(s, IMG("activity_geofence.png"), Inches(0.55), Inches(2.05), Inches(6.0), Inches(4.85), border=True)
numbered_steps(s, Inches(6.9), Inches(2.05), Inches(5.85), [
    "Phụ huynh vẽ vùng an toàn trên bản đồ — hình tròn hoặc hành lang đường đi",
    "Backend lưu cấu hình, áp dụng ngay cho các bản tin vị trí tiếp theo",
    "Mỗi tọa độ mới nhận được, backend kiểm tra có nằm trong vùng an toàn không",
    "Nếu vi phạm: tạo bản ghi cảnh báo, gửi Telegram/Email và đẩy qua WebSocket",
    "Ứng dụng hiển thị cảnh báo kèm tên vùng và vị trí vi phạm",
], size=13.5, item_h=0.62, gap=0.12, badge_color=CORAL)
rounded_rect(s, Inches(6.9), Inches(6.55), Inches(5.85), Inches(0.65), CORAL_LT, radius=0.12)
text(s, Inches(7.1), Inches(6.63), Inches(5.5), Inches(0.5),
     "Thời gian chờ tối thiểu 60 giây giữa 2 cảnh báo liên tiếp của cùng một vùng — chống spam",
     size=11.5, italic=True, color=CORAL, line_spacing=1.1)

# =====================================================================
# SLIDE 5D — LUỒNG 3: XỬ LÝ TÍN HIỆU SOS
# =====================================================================
s = content_slide(SECTION_TAGS["design"], "Luồng nghiệp vụ 3: Xử lý tín hiệu SOS khẩn cấp",
                   subtitle="Tác nhân: Thiết bị IoT · Backend · Dịch vụ thông báo · Phụ huynh")
image_fit(s, IMG("activity_sos.png"), Inches(0.55), Inches(2.05), Inches(6.0), Inches(4.85), border=True)
numbered_steps(s, Inches(6.9), Inches(2.05), Inches(5.85), [
    "Trẻ nhấn nút SOS vật lý trên thiết bị",
    "Thiết bị đọc tọa độ hiện tại, gửi ngay tín hiệu SOS qua MQTT (ưu tiên cao)",
    "Backend nhận tín hiệu, tạo bản ghi cảnh báo khẩn cấp",
    "Gửi đồng thời (asyncio.gather) qua Telegram kèm link Google Maps và qua WebSocket",
    "Ứng dụng hiển thị popup cảnh báo khẩn cùng vị trí hiện tại trên bản đồ",
], size=13.5, item_h=0.62, gap=0.12, badge_color=CORAL)
rounded_rect(s, Inches(6.9), Inches(6.55), Inches(5.85), Inches(0.65), CORAL_LT, radius=0.12)
text(s, Inches(7.1), Inches(6.63), Inches(5.5), Inches(0.5),
     "Nếu mất kết nối MQTT: thiết bị lưu sự kiện SOS và gửi lại ngay khi kết nối được khôi phục",
     size=11.5, italic=True, color=CORAL, line_spacing=1.1)

# =====================================================================
# SLIDE 6 — THIẾT KẾ CSDL
# =====================================================================
s = content_slide(SECTION_TAGS["design"], "Thiết kế cơ sở dữ liệu (MongoDB)")
image_fit(s, IMG("er_diagram.png"), Inches(6.55), Inches(1.75), Inches(6.2), Inches(5.1))
bullets(s, Inches(0.55), Inches(1.8), Inches(5.7), Inches(4.7), [
    ("users — tài khoản phụ huynh", 0),
    ("device_permissions — quyền sở hữu/chia sẻ thiết bị (1–N với users)", 0),
    ("device_configs — cấu hình thiết bị (1–1 với thiết bị)", 0),
    ("locations — chuỗi tọa độ GPS theo thời gian (1–N theo thiết bị)", 0),
    ("geofence_configs — cấu hình vùng an toàn (mảng geofences: mode, tâm, path)", 0),
    ("alerts — lịch sử cảnh báo (1–N theo thiết bị)", 0),
], size=13.5, space_after=10)
rounded_rect(s, Inches(0.55), Inches(6.05), Inches(5.7), Inches(0.95), DARK_BOX, radius=0.08)
text(s, Inches(0.78), Inches(6.16), Inches(5.3), Inches(0.75),
     "Chỉ mục phức hợp (deviceId, timestamp) + TTL trên receivedAt\n→ tự động xóa dữ liệu cũ sau 30 ngày",
     size=12.5, color=WHITE, line_spacing=1.2)

# =====================================================================
# SLIDE 7 — CÔNG NGHỆ PHẦN CỨNG
# =====================================================================
s = content_slide(SECTION_TAGS["tech"], "Công nghệ phần cứng: ESP32 + SIM7600")
card(s, Inches(0.55), Inches(1.8), Inches(5.95), Inches(4.7), "ESP32", [
    "Dual-core Xtensa LX6, xung nhịp 240 MHz",
    "Wi-Fi 802.11 b/g/n, Bluetooth/BLE tích hợp sẵn",
    "UART, SPI, I²C, ADC — giao tiếp với SIM7600 qua UART2",
    "Ngắt ngoài (GPIO) xử lý sự kiện nhấn nút SOS tức thì",
], body_size=15, title_size=19)
card(s, Inches(6.8), Inches(1.8), Inches(5.95), Inches(4.7), "SIM7600", [
    "Tích hợp GNSS (GPS/GLONASS/BeiDou) + modem 4G LTE Cat-1 trong 1 module",
    "Độ chính xác GPS ~2,5 m CEP ngoài trời",
    "Giao tiếp qua UART bằng AT command, trừu tượng hoá qua thư viện TinyGSM",
    "1 anten duy nhất phục vụ cả định vị và kết nối di động",
], body_size=15, title_size=19, title_color=CORAL, bg=CORAL_LT)
text(s, Inches(0.55), Inches(6.65), Inches(12.2), Inches(0.4),
     "→ Một module duy nhất thay thế 2 module GPS + modem rời riêng biệt",
     size=14, bold=True, color=TEXT_DARK, align=PP_ALIGN.CENTER)

# =====================================================================
# SLIDE 8 — GIAO THỨC & BACKEND
# =====================================================================
s = content_slide(SECTION_TAGS["tech"], "Giao thức truyền thông & Backend")
card(s, Inches(0.55), Inches(1.8), Inches(3.95), Inches(4.7), "MQTT", [
    "Giao thức nhẹ, mô hình publish/subscribe",
    "QoS 1 — đảm bảo dữ liệu không mất trên mạng di động không ổn định",
    "Broker: EMQX public (broker.emqx.io)",
], body_size=13.5, title_size=17)
card(s, Inches(4.7), Inches(1.8), Inches(3.95), Inches(4.7), "FastAPI", [
    "Framework Python bất đồng bộ (async/await)",
    "Xử lý đồng thời MQTT callback, WebSocket, REST API, MongoDB",
    "Tự sinh tài liệu API theo chuẩn OpenAPI",
], body_size=13.5, title_size=17)
card(s, Inches(8.85), Inches(1.8), Inches(3.95), Inches(4.7), "MongoDB", [
    "NoSQL hướng tài liệu, lưu JSON linh hoạt",
    "Time Series Collection cho dữ liệu chuỗi thời gian",
    "Chỉ mục địa lý 2dsphere cho lịch sử hành trình",
], body_size=13.5, title_size=17, title_color=CORAL, bg=CORAL_LT)

# =====================================================================
# SLIDE 9 — ỨNG DỤNG DI ĐỘNG
# =====================================================================
s = content_slide(SECTION_TAGS["tech"], "Ứng dụng di động: Flutter + OpenStreetMap")
bullets(s, Inches(0.55), Inches(1.8), Inches(7.2), Inches(4.9), [
    ("Flutter (Dart) — một codebase cho cả Android và iOS", 0),
    ("Biên dịch AOT sang mã máy, engine đồ hoạ riêng (Skia/Impeller) → 60fps nhất quán", 1),
    ("flutter_map + OpenStreetMap — bản đồ tương tác, miễn phí, không cần API key", 0),
    ("web_socket_channel — nhận cập nhật vị trí thời gian thực từ backend", 0),
    ("Provider pattern — quản lý trạng thái ứng dụng", 0),
], size=16, space_after=14)
image_fit(s, IMG("tk_giaodien.png"), Inches(7.95), Inches(1.8), Inches(4.85), Inches(4.9))

# =====================================================================
# SLIDE 10 — ĐÓNG GÓP 1
# =====================================================================
s = content_slide(SECTION_TAGS["contrib"], "Đóng góp 1: Tích hợp GPS + kết nối di động trên một module")
card(s, Inches(0.55), Inches(1.8), Inches(5.95), Inches(2.55), "Bài toán", [
    "2 module riêng biệt (GPS + modem) → tăng linh kiện, cần 2 anten",
    "Firmware phải điều phối song song 2 luồng UART khác nhau → phức tạp, dễ xung đột tài nguyên",
], body_size=14, title_size=16, bg=CORAL_LT, title_color=CORAL)
card(s, Inches(6.8), Inches(1.8), Inches(5.95), Inches(2.55), "Giải pháp", [
    "SIM7600: tích hợp GNSS + 4G LTE trong 1 linh kiện, 1 anten, 1 UART",
    "Firmware theo state machine: INIT → GSM_CONNECT → MQTT_CONNECT → GPS_WAIT_FIX → TRACKING",
    "Tự động quay lại GSM_CONNECT khi mất mạng, không gián đoạn phiên GPS",
], body_size=14, title_size=16)
rounded_rect(s, Inches(0.55), Inches(4.55), Inches(12.2), Inches(2.05), DARK_BOX, radius=0.06)
text(s, Inches(0.85), Inches(4.72), Inches(11.6), Inches(0.35), "Kết quả đạt được", size=15, bold=True, color=WHITE)
bullets(s, Inches(0.85), Inches(5.12), Inches(11.6), Inches(1.35), [
    "Giảm số linh kiện chính trên bo mạch, đơn giản hoá sơ đồ nguyên lý và bố cục PCB",
    "Firmware chỉ quản lý 1 giao diện UART, giảm nguy cơ tranh chấp tài nguyên",
    "Thiết bị tự phục hồi sau khi mất tín hiệu GPS hoặc rớt mạng mà không cần can thiệp thủ công",
], size=13.5, color=WHITE, bullet_color=WHITE, space_after=6)

# =====================================================================
# SLIDE 10B — LƯU ĐỒ FIRMWARE ESP32
# =====================================================================
s = content_slide(SECTION_TAGS["contrib"], "Đóng góp 1 (tiếp): Lưu đồ firmware ESP32",
                   subtitle="Ba giai đoạn vận hành: khởi tạo → vòng lặp chính → xử lý ngắt SOS")
fw_items = [
    ("firmware_setup.png", "Khởi tạo (Setup)",
     "UART → reset/bật GPS SIM7600 → lặp đến khi có GPRS + MQTT → đăng ký ngắt SOS (GPIO0)"),
    ("firmware_loop.png", "Vòng lặp chính (Main Loop)",
     "Duy trì MQTT, kiểm tra cờ sosPending, mỗi 5s đọc GPS và publish nếu đã có định vị"),
    ("firmware_sos.png", "Xử lý ngắt SOS",
     "Khử rung 50ms → đọc GPS (fallback vị trí cache) → cooldown 10s → publish SOS + nháy LED"),
]
fw_w = Inches(3.95)
fw_gap = Inches(0.2)
fw_startx = Inches(0.55)
for i, (fname, caption, desc) in enumerate(fw_items):
    x = fw_startx + i * (fw_w + fw_gap)
    image_fit(s, IMG(fname), x, Inches(1.85), fw_w - Inches(0.3), Inches(4.15))
    text(s, x, Inches(6.05), fw_w, Inches(0.35), caption, size=14.5, bold=True, color=TEXT_DARK,
         align=PP_ALIGN.CENTER)
    text(s, x, Inches(6.4), fw_w, Inches(0.75), desc, size=10.5, color=TEXT_MUTE,
         align=PP_ALIGN.CENTER, line_spacing=1.15)

# =====================================================================
# SLIDE 11 — ĐÓNG GÓP 2
# =====================================================================
s = content_slide(SECTION_TAGS["contrib"], "Đóng góp 2: Geofencing — chế độ tâm × hình dạng vùng")
text(s, Inches(0.55), Inches(1.62), Inches(12.2), Inches(0.35),
     "Hai trục lựa chọn độc lập:  (1) chế độ tâm — Fixed / Mobile   ·   (2) hình dạng vùng (chỉ áp dụng khi Fixed) — Tròn / Đường đi",
     size=12.5, italic=True, color=TEXT_MUTE)

# LEFT: khối lớn "FIXED" chứa 2 hình dạng con
rounded_rect(s, Inches(0.55), Inches(2.05), Inches(7.3), Inches(4.55), TEAL_LT, radius=0.04)
rect(s, Inches(0.55), Inches(2.05), Inches(0.07), Inches(4.55), TEAL)
text(s, Inches(0.85), Inches(2.18), Inches(6.8), Inches(0.4),
     "Chế độ FIXED — tâm cố định do phụ huynh đặt", size=16, bold=True, color=TEXT_DARK)

rounded_rect(s, Inches(0.85), Inches(2.68), Inches(6.7), Inches(1.75), WHITE, radius=0.06, line=True)
text(s, Inches(1.05), Inches(2.78), Inches(6.3), Inches(0.35), "a. Hình tròn  (path rỗng)",
     size=14, bold=True, color=TEAL)
bullets(s, Inches(1.05), Inches(3.14), Inches(6.3), Inches(1.2), [
    "Khoảng cách Haversine từ vị trí hiện tại đến tâm cố định",
    "Vi phạm khi khoảng cách d > radiusM",
], size=12.5, space_after=4, bullet_color=TEAL)

rounded_rect(s, Inches(0.85), Inches(4.55), Inches(6.7), Inches(1.9), WHITE, radius=0.06, line=True)
text(s, Inches(1.05), Inches(4.65), Inches(6.3), Inches(0.35), "b. Hành lang đường đi  (path ≥ 2 điểm)",
     size=14, bold=True, color=TEAL)
bullets(s, Inches(1.05), Inches(5.01), Inches(6.3), Inches(1.35), [
    "Chiếu vị trí lên từng đoạn thẳng của tuyến đường (point-to-segment)",
    "Vi phạm khi khoảng cách nhỏ nhất trong các đoạn > bán kính hành lang",
], size=12.5, space_after=4, bullet_color=TEAL)

# RIGHT: khối "MOBILE" — tâm động theo GPS điện thoại
card(s, Inches(8.1), Inches(2.05), Inches(4.65), Inches(4.55), "Chế độ MOBILE — tâm động", [
    "Tâm KHÔNG cố định: đồng bộ theo vị trí GPS điện thoại phụ huynh mỗi ~10 giây",
    "Luôn là hình tròn quanh vị trí tức thời — không áp dụng hình dạng đường đi",
    "Phù hợp khi muốn 'trẻ luôn cách tôi trong bán kính X mét' thay vì gắn với một điểm cố định",
], body_size=13.5, title_size=16, title_color=CORAL, bg=CORAL_LT)

text(s, Inches(0.55), Inches(6.75), Inches(12.2), Inches(0.4),
     "→ Trường mode lưu 'fixed' / 'mobile'; hình dạng tròn hay đường đi được suy ra từ dữ liệu path, độc lập với mode",
     size=12.5, bold=True, color=TEXT_DARK, align=PP_ALIGN.CENTER)

# =====================================================================
# SLIDE 12 — ĐÓNG GÓP 3
# =====================================================================
s = content_slide(SECTION_TAGS["contrib"], "Đóng góp 3: Kiến trúc lai MQTT – WebSocket")
bullets(s, Inches(0.55), Inches(1.8), Inches(6.4), Inches(3.0), [
    "Không giao thức đơn lẻ nào tối ưu cho cả 2 chặng: thiết bị→server (MQTT) và server→ứng dụng (WebSocket)",
    "Backend đóng vai trò cầu nối (bridge) giữa hai giao thức",
], size=15, space_after=10)
card(s, Inches(0.55), Inches(3.4), Inches(6.4), Inches(3.1), "Cơ chế cầu nối", [
    "MQTT callback (đồng bộ, thread riêng) → đẩy vào asyncio.Queue",
    "Coroutine nền đọc hàng đợi → broadcast qua WebSocket đến các app đang theo dõi",
    "Tách bạch: thiết bị chỉ biết MQTT broker, app chỉ biết WebSocket backend",
], body_size=13.5, title_size=16)
image_fit(s, IMG("sequence_vitri.png"), Inches(7.2), Inches(1.8), Inches(5.6), Inches(4.7))
rounded_rect(s, Inches(0.55), Inches(6.65), Inches(3.1), Inches(0.65), DARK_BOX_2, radius=0.15)
text(s, Inches(0.55), Inches(6.65), Inches(3.1), Inches(0.65), "< 2s   độ trễ end-to-end",
     size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
rounded_rect(s, Inches(3.85), Inches(6.65), Inches(3.1), Inches(0.65), DARK_BOX_2, radius=0.15)
text(s, Inches(3.85), Inches(6.65), Inches(3.1), Inches(0.65), "< 300ms   độ trễ truyền thông",
     size=12.5, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

# =====================================================================
# SLIDE 12B — SD02: TUẦN TỰ CẢNH BÁO VI PHẠM GEOFENCE
# =====================================================================
s = content_slide(SECTION_TAGS["contrib"], "Biểu đồ tuần tự (SD02): Cảnh báo vi phạm vùng an toàn")
image_fit(s, IMG("sequence_geofence.png"), Inches(0.55), Inches(1.75), Inches(12.25), Inches(4.2))
bullets(s, Inches(0.8), Inches(6.05), Inches(11.7), Inches(1.05), [
    "Tiếp nối luồng SD01 (cập nhật vị trí): Backend kiểm tra geofence ngay khi nhận tọa độ mới",
    "Khi phát hiện vi phạm: lưu bản ghi cảnh báo → gửi Telegram/Email → phát cảnh báo qua WebSocket tới ứng dụng",
], size=13.5, space_after=5)

# =====================================================================
# SLIDE 12C — SD03: TUẦN TỰ XỬ LÝ TÍN HIỆU SOS
# =====================================================================
s = content_slide(SECTION_TAGS["contrib"], "Biểu đồ tuần tự (SD03): Gửi tín hiệu SOS khẩn cấp")
image_fit(s, IMG("sequence_sos.png"), Inches(0.55), Inches(1.75), Inches(12.25), Inches(4.2))
bullets(s, Inches(0.8), Inches(6.05), Inches(11.7), Inches(1.05), [
    "ESP32 đọc tọa độ hiện tại và gửi ngay tín hiệu SOS qua MQTT khi trẻ nhấn nút",
    "Backend phát hiện, lưu cảnh báo khẩn và gửi ngay thông báo SOS kèm liên kết Google Maps tới phụ huynh qua Telegram",
], size=13.5, space_after=5)

# =====================================================================
# SLIDE 13 — ĐÓNG GÓP 4
# =====================================================================
s = content_slide(SECTION_TAGS["contrib"], "Đóng góp 4: Cảnh báo đa kênh song song")
chan_w = Inches(3.95)
channels = [
    ("WebSocket", "Trong ứng dụng, độ trễ thấp nhất; chỉ hoạt động khi app đang mở", WHITE),
    ("Telegram Bot", "Kèm tọa độ + link Google Maps; hoạt động độc lập với app", WHITE),
    ("Email", "Lớp dự phòng cuối cùng, lưu lại bằng chứng để tra cứu sau", WHITE),
]
for i, (name, desc, color) in enumerate(channels):
    x = Inches(0.55) + i * (chan_w + Inches(0.2))
    rounded_rect(s, x, Inches(1.85), chan_w, Inches(2.5), DARK_BOX, radius=0.07)
    text(s, x + Inches(0.25), Inches(2.05), chan_w - Inches(0.5), Inches(0.45), name,
         size=17, bold=True, color=color)
    text(s, x + Inches(0.25), Inches(2.6), chan_w - Inches(0.5), Inches(1.6), desc,
         size=12.5, color=WHITE, line_spacing=1.2)
bullets(s, Inches(0.55), Inches(4.65), Inches(12.2), Inches(1.4), [
    "Cả 3 kênh gọi song song bằng asyncio.gather — không chờ tuần tự từng kênh",
    "Mọi sự kiện cảnh báo đều ghi vào collection alerts để tra cứu lịch sử sau này",
], size=15, space_after=8)
rounded_rect(s, Inches(0.55), Inches(6.15), Inches(12.2), Inches(0.75), TEAL_LT, radius=0.1)
text(s, Inches(0.85), Inches(6.32), Inches(11.6), Inches(0.4),
     "Kết quả: cả 3 kênh gửi thành công trong < 3 giây — Telegram nhanh nhất (< 1 giây)",
     size=14, bold=True, color=TEAL, align=PP_ALIGN.CENTER)

# =====================================================================
# SLIDE 14 — DEMO GIAO DIỆN
# =====================================================================
s = content_slide(SECTION_TAGS["result"], "Demo giao diện ứng dụng")
demo_items = [
    ("demo_map.png", "Theo dõi vị trí thời gian thực"),
    ("demo_geofence.png", "Tạo & chỉnh sửa vùng an toàn"),
    ("demo_alert.png", "Cảnh báo vi phạm / SOS"),
    ("demo_stats.png", "Thống kê hành trình"),
]
# Luoi 2x2: anh demo co ty le rong (2.22:1) nen xep 1 hang 4 anh se rat be
# va thua nhieu khoang trang tren/duoi; xep 2x2 giup anh to hon ro ret.
demo_col_w = Inches(5.9)
demo_gap_x = Inches(0.43)
demo_img_h = Inches(2.15)
demo_row_y = [Inches(1.7), Inches(4.42)]
for i, (fname, cap) in enumerate(demo_items):
    col, row = i % 2, i // 2
    x = Inches(0.55) + col * (demo_col_w + demo_gap_x)
    y = demo_row_y[row]
    image_fit(s, IMG(fname), x, y, demo_col_w, demo_img_h, border=True)
    text(s, x, y + demo_img_h + Inches(0.1), demo_col_w, Inches(0.34), cap,
         size=14, bold=True, color=TEXT_DARK, align=PP_ALIGN.CENTER, line_spacing=1.05)

# =====================================================================
# SLIDE 15 — KẾT QUẢ KIỂM THỬ & SỐ LIỆU
# =====================================================================
s = content_slide(SECTION_TAGS["result"], "Kết quả kiểm thử & số liệu đo thực tế")
mvals = [("< 2s", "Độ trễ end-to-end"), ("~3,2 m", "Độ chính xác GPS"), ("5s", "Chu kỳ cập nhật"), ("10/10", "Test case đạt")]
mw = Inches(2.95)
gap = Inches(0.2)
for i, (val, lab) in enumerate(mvals):
    x = Inches(0.55) + i * (mw + gap)
    metric_card(s, x, Inches(1.8), mw, Inches(1.3), val, lab)
data2 = [
    ["Mã", "Mô tả", "Kết quả mong đợi", "KQ"],
    ["TC01", "Thiết bị gửi tọa độ hợp lệ", "Bản đồ cập nhật < 3 giây", "Đạt"],
    ["TC03", "Trẻ ra ngoài vùng an toàn", "Sinh cảnh báo, gửi Telegram/Email", "Đạt"],
    ["TC05", "Hai cảnh báo liên tiếp trong cooldown", "Chỉ gửi một cảnh báo", "Đạt"],
    ["TC07", "Nhấn nút SOS", "Gửi cảnh báo kèm link bản đồ qua Telegram", "Đạt"],
    ["TC10", "Mất kết nối WebSocket", "Ứng dụng tự kết nối lại", "Đạt"],
]
make_table(s, Inches(0.55), Inches(3.4), Inches(12.25), Inches(3.15), data2,
           col_w_ratio=[0.8, 3.6, 3.6, 0.9], highlight_col=3, header_size=13, body_size=12.5)
text(s, Inches(0.55), Inches(6.72), Inches(12.2), Inches(0.35),
     "Kiểm thử hộp đen (black-box) + tích hợp đầu cuối, mô phỏng publish MQTT không cần thiết bị vật lý",
     size=11.5, italic=True, color=TEXT_MUTE)

# =====================================================================
# SLIDE 16 — SO SÁNH TỔNG KẾT & TRIỂN KHAI
# =====================================================================
s = content_slide(SECTION_TAGS["result"], "So sánh tổng kết & triển khai thực tế")
card(s, Inches(0.55), Inches(1.8), Inches(5.95), Inches(3.4), "So với sản phẩm thị trường", [
    "Đề tài đáp ứng đầy đủ 5/5 tiêu chí: độc lập, thời gian thực, geofence linh hoạt, cảnh báo đa kênh, chi phí thấp",
    "Mỗi sản phẩm thương mại (Life360, AirTag, GPS Tracker) chỉ đáp ứng được một phần",
], body_size=15.5, title_size=18)
card(s, Inches(6.8), Inches(1.8), Inches(5.95), Inches(3.4), "Triển khai thực tế", [
    "Backend: đóng gói Docker, chạy trên Render (Uvicorn/ASGI)",
    "MongoDB Atlas — gói miễn phí, sao lưu tự động",
    "EMQX public broker — không cần tự vận hành hạ tầng",
    "Ứng dụng: đóng gói APK cài trực tiếp lên máy phụ huynh",
], body_size=15.5, title_size=18, title_color=CORAL, bg=CORAL_LT)
rounded_rect(s, Inches(0.55), Inches(5.45), Inches(12.2), Inches(1.35), DARK_BOX, radius=0.08)
text(s, Inches(0.85), Inches(5.62), Inches(11.6), Inches(0.35), "Điểm nổi bật", size=15, bold=True, color=WHITE)
text(s, Inches(0.85), Inches(6.0), Inches(11.6), Inches(0.7),
     "Chi phí vận hành duy nhất: gói data SIM di động + phí hosting backend — toàn bộ hạ tầng nền (MQTT broker, database, bản đồ) đều dùng dịch vụ miễn phí",
     size=14, color=WHITE, line_spacing=1.25)

# =====================================================================
# SLIDE 17 — KẾT LUẬN
# =====================================================================
s = content_slide(SECTION_TAGS["concl"], "Kết luận")
card(s, Inches(0.55), Inches(2.0), Inches(5.95), Inches(3.6), "Kết quả đạt được", [
    "Xây dựng hoàn chỉnh hệ thống 3 thành phần hoạt động đồng bộ",
    "Thiết bị hoạt động độc lập, không cần điện thoại của trẻ",
    "Theo dõi thời gian thực với độ trễ < 2 giây",
    "Geofencing linh hoạt 2 chế độ, cảnh báo đa kênh",
    "Chi phí vận hành thấp nhờ dịch vụ miễn phí (EMQX, MongoDB Atlas, OSM, Telegram)",
], body_size=14.5, title_size=18)
card(s, Inches(6.8), Inches(2.0), Inches(5.95), Inches(3.6), "Hạn chế còn tồn tại", [
    "Chưa có thiết kế PCB tùy chỉnh — dùng module phát triển rời",
    "Xác thực JWT chưa có refresh token / thu hồi token",
    "Kênh MQTT chưa được mã hoá TLS",
    "Tính năng chia sẻ thiết bị đa phụ huynh: có DB, chưa hoàn thiện giao diện",
], body_size=14.5, title_size=18, title_color=CORAL, bg=CORAL_LT)

# =====================================================================
# SLIDE 18 — HƯỚNG PHÁT TRIỂN
# =====================================================================
s = content_slide(SECTION_TAGS["concl"], "Hướng phát triển")
card(s, Inches(0.55), Inches(2.1), Inches(5.95), Inches(3.5), "Phần cứng & Firmware", [
    "Thiết kế PCB tùy chỉnh và đóng gói sản phẩm nhỏ gọn, đeo được",
    "Chế độ ngủ sâu (deep sleep) & module tiết kiệm điện NB-IoT/LTE-M",
    "Mở rộng đa thiết bị và phân quyền chia sẻ theo dõi",
], body_size=16, title_size=18)
card(s, Inches(6.8), Inches(2.1), Inches(5.95), Inches(3.5), "Phần mềm & Bảo mật", [
    "Phân tích hành vi bất thường từ dữ liệu chuỗi thời gian",
    "Chế độ hoạt động ngoại tuyến — lưu flash khi mất kết nối",
    "Bảo mật nâng cao: refresh token, TLS cho MQTT, xác thực hai yếu tố (2FA)",
], body_size=16, title_size=18, title_color=CORAL, bg=CORAL_LT)

# =====================================================================
# SLIDE 19 — CẢM ƠN
# =====================================================================
s = new_slide()
set_bg(s, NAVY)
rect(s, 0, 0, SLIDE_W, Inches(0.12), GOLD)
rect(s, 0, Inches(7.38), SLIDE_W, Inches(0.12), GOLD)
end_logo_h = Inches(0.75)
end_logo_w = Emu(int(end_logo_h * LOGO_AR))
s.shapes.add_picture(LOGO, (SLIDE_W - end_logo_w) / 2, Inches(0.55), height=end_logo_h, width=end_logo_w)
text(s, Inches(1), Inches(2.75), Inches(11.3), Inches(1.2), "XIN CẢM ƠN",
     size=48, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
rect(s, Inches(5.67), Inches(4.0), Inches(2), Inches(0.035), GOLD)
text(s, Inches(1.5), Inches(4.2), Inches(10.3), Inches(0.5), "Câu hỏi & Thảo luận",
     size=18, color=GOLD, align=PP_ALIGN.CENTER)
text(s, Inches(1), Inches(6.6), Inches(11.3), Inches(0.4),
     "Nguyễn Đức Dương  ·  MSSV 20225122  ·  SOICT – Đại học Bách khoa Hà Nội",
     size=12, color=RGBColor(0x8C, 0xA3, 0xB3), align=PP_ALIGN.CENTER)

prs.save(OUT_PATH)
print("Saved:", OUT_PATH)
print("Total slides:", len(prs.slides))
